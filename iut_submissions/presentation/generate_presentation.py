"""
Thesis Pre-Defense Presentation Generator
==========================================
Generates a complete .pptx file for the thesis pre-defense presentation:
"Orthogonal Defect Classification on Defects4J Using an LLM-Driven Scientific Approach"

Usage: python generate_presentation.py
Output: Thesis_PreDefense_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# DESIGN CONSTANTS
# ============================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette – Professional Navy + Teal
PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)       # Deep Navy
ACCENT1 = RGBColor(0x2E, 0x86, 0xAB)       # Teal
ACCENT2 = RGBColor(0x48, 0xA9, 0xA6)       # Seafoam
ACCENT3 = RGBColor(0xE8, 0x6F, 0x51)       # Coral
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG    = RGBColor(0xF5, 0xF7, 0xFA)   # Very light blue-gray
TEXT_DARK   = RGBColor(0x2C, 0x3E, 0x50)
TEXT_MED    = RGBColor(0x5D, 0x6D, 0x7E)
TEXT_LIGHT  = RGBColor(0x85, 0x92, 0x9E)
BORDER      = RGBColor(0xD5, 0xDB, 0xE1)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)

FONT = "Calibri"

# ============================================================
# HELPER FUNCTIONS
# ============================================================

_page = [0]   # mutable page counter


def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _page_num(slide, num):
    tb = slide.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT


def _bot_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.3),
                               SLIDE_WIDTH, Inches(0.2))
    _solid(s, PRIMARY)
    _no_line(s)


def _accent_line(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15),
                               Inches(2.0), Inches(0.04))
    _solid(s, ACCENT1)
    _no_line(s)


def _section_label(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(9)
    p.font.color.rgb = ACCENT1
    p.font.name = FONT
    p.font.bold = True


def _slide_title(slide, title, section=None):
    if section:
        _section_label(slide, section)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(11.5), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY
    p.font.name = FONT
    p.font.bold = True
    _accent_line(slide)


def new_slide(prs, title, section=None):
    """Create a standard content slide and return it."""
    _page[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    _slide_title(slide, title, section)
    _bot_bar(slide)
    if _page[0] > 1:
        _page_num(slide, _page[0])
    return slide


def _txt(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _box(slide, left, top, w, h, fill_c=LIGHT_BG, line_c=BORDER, line_w=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(left), Inches(top), Inches(w), Inches(h))
    _solid(s, fill_c)
    s.line.color.rgb = line_c
    s.line.width = Pt(line_w)
    return s


def _pill(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(left), Inches(top), Inches(w), Inches(h))
    _solid(s, color)
    _no_line(s)
    return s


def _p(tf, text, size=12, color=TEXT_DARK, bold=False, italic=False,
       align=None, after=Pt(4), first=False):
    """Append (or reuse first) paragraph."""
    if first:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.bold = bold
    p.font.italic = italic
    if align:
        p.alignment = align
    p.space_after = after
    return p


# ============================================================
# MAIN GENERATOR
# ============================================================

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ==================================================================
    # SLIDE 1 — TITLE
    # ==================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _page[0] = 1
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = PRIMARY

    # top accent bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               SLIDE_WIDTH, Inches(0.08))
    _solid(s, ACCENT1); _no_line(s)

    # university & lab
    tf = _txt(slide, 1, 0.5, 11.333, 0.6)
    _p(tf, "ISLAMIC UNIVERSITY OF TECHNOLOGY", 13, ACCENT2, True,
       align=PP_ALIGN.CENTER, after=Pt(4), first=True)
    _p(tf, "Network and Data Analysis Group (NDAG)", 11,
       RGBColor(0xA0, 0xC4, 0xD0), align=PP_ALIGN.CENTER)

    # divider
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.4),
                               Inches(5.333), Inches(0.03))
    _solid(s, ACCENT1); _no_line(s)

    # thesis title
    tf = _txt(slide, 0.8, 1.7, 11.733, 1.6)
    _p(tf, "Orthogonal Defect Classification on Defects4J", 34, WHITE, True,
       align=PP_ALIGN.CENTER, after=Pt(4), first=True)
    _p(tf, "Using an LLM-Driven Scientific Approach", 30, ACCENT2, True,
       align=PP_ALIGN.CENTER)

    # subtitle
    tf = _txt(slide, 1, 3.5, 11.333, 0.4)
    _p(tf, "Thesis Pre-Defense Presentation", 14,
       RGBColor(0xA0, 0xC4, 0xD0), align=PP_ALIGN.CENTER, first=True)

    # divider 2
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.1),
                               Inches(2.333), Inches(0.02))
    _solid(s, ACCENT2); _no_line(s)

    # team members
    tf = _txt(slide, 1, 4.4, 11.333, 1.5)
    _p(tf, "Presented by", 10, TEXT_LIGHT, align=PP_ALIGN.CENTER, after=Pt(6), first=True)
    for m in ["Hasin Mahtab Alvee  (210042174)",
              "Mohammad Nahiyan Kabir  (210042168)",
              "Md. Sakib Hossain  (210042133)"]:
        _p(tf, m, 13, WHITE, align=PP_ALIGN.CENTER, after=Pt(2))

    # supervisors
    tf = _txt(slide, 1, 6.0, 11.333, 1.2)
    _p(tf, "Supervised by", 10, TEXT_LIGHT, align=PP_ALIGN.CENTER, after=Pt(4), first=True)
    _p(tf, "Lutfun Nahar Lota, Assistant Professor  |  Ishmam Tashdeed, Lecturer",
       13, ACCENT2, align=PP_ALIGN.CENTER, after=Pt(2))
    _p(tf, "Signal and Speech Processing Laboratory (SSL)", 11,
       RGBColor(0xA0, 0xC4, 0xD0), align=PP_ALIGN.CENTER)

    # bottom bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35),
                               SLIDE_WIDTH, Inches(0.15))
    _solid(s, ACCENT1); _no_line(s)

    # ==================================================================
    # SLIDE 2 — OUTLINE
    # ==================================================================
    slide = new_slide(prs, "Presentation Outline", "OVERVIEW")

    left_items = [
        ("01", "Introduction", "Defects4J & Bug Categorization"),
        ("02", "Literature Review", "Existing Classifications & LLM Approaches"),
        ("03", "Problem Statement", "Research Gap & Motivation"),
        ("04", "Objectives", "Aims of the Study"),
        ("05", "Proposed Methodology", "Pipeline Architecture & Design"),
    ]
    right_items = [
        ("06", "Results / Findings", "Classification \u0026 Comparison Data"),
        ("07", "Discussion / Analysis", "Divergence Patterns \u0026 Case Study"),
        ("08", "Limitations", "Known Constraints"),
        ("09", "Future Work", "Roadmap \u0026 Extensions"),
        ("10", "Conclusion", "Summary \u0026 Closing"),
    ]

    for items, x_base, circle_clr in [
        (left_items, 0.8, PRIMARY),
        (right_items, 7.0, ACCENT1),
    ]:
        for i, (num, title, desc) in enumerate(items):
            y = 1.6 + i * 1.05
            # circle
            s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_base), Inches(y),
                                       Inches(0.5), Inches(0.5))
            _solid(s, circle_clr); _no_line(s)
            s.text_frame.word_wrap = False
            p = s.text_frame.paragraphs[0]
            p.text = num; p.font.size = Pt(12); p.font.color.rgb = WHITE
            p.font.name = FONT; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
            s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            tx = x_base + 0.7
            tf = _txt(slide, tx, y, 4.5, 0.3)
            _p(tf, title, 16, PRIMARY, True, first=True)
            tf2 = _txt(slide, tx, y + 0.28, 4.5, 0.25)
            _p(tf2, desc, 11, TEXT_MED, first=True)

    # ==================================================================
    # SLIDE 3 — INTRODUCTION: DEFECTS4J  (SP1)
    # ==================================================================
    slide = new_slide(prs, "Introduction to Defects4J", "INTRODUCTION")

    tf = _txt(slide, 0.6, 1.4, 6.5, 5.0)
    _p(tf, "What is Defects4J?", 18, PRIMARY, True, after=Pt(8), first=True)
    for b in [
        "• A curated database of real, reproducible Java bugs with developer-written patches [1]",
        "• Provides per-bug artifacts: bug reports, triggering tests, buggy code, fixed code, diffs",
        "• Standard benchmark for empirical SE research (fault localization, APR, testing)",
        "• Grown to 864+ active bugs across 17 Java projects [3]",
        "• Includes CLI tools for checkout, compilation, testing, and coverage analysis",
    ]:
        _p(tf, b, 13, TEXT_DARK)

    _p(tf, "", 6)
    _p(tf, "Continued Relevance", 18, PRIMARY, True, after=Pt(8))
    for b in [
        "• Actively used in 2024 fault localization studies [3]",
        "• Standard benchmark for large-scale APR experiments [5]",
        "• Manual inspections confirm diagnostic value of test-suite artifacts [6]",
    ]:
        _p(tf, b, 13, TEXT_DARK)

    # right info box
    s = _box(slide, 7.6, 1.4, 5.0, 4.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    _p(tf, "Key Artifacts per Bug", 16, PRIMARY, True,
       align=PP_ALIGN.CENTER, after=Pt(10), first=True)
    for icon, t, d in [
        ("\U0001F4CB", "Bug Report", "Original issue description"),
        ("\U0001F9EA", "Triggering Tests", "Tests that expose the bug"),
        ("\U0001F4BB", "Buggy Source Code", "Pre-fix Java source"),
        ("\u2705", "Fixed Source Code", "Developer-patched code"),
        ("\U0001F4CA", "Code Diffs", "Buggy-to-Fixed changes"),
        ("\u2699\uFE0F", "CLI Tools", "Checkout, compile, test, coverage"),
    ]:
        _p(tf, f"{icon}  {t}", 13, TEXT_DARK, True, after=Pt(0))
        _p(tf, f"     {d}", 11, TEXT_MED, after=Pt(8))

    # nuance callout
    s = _box(slide, 0.6, 6.2, 12.0, 0.8,
             RGBColor(0xFF, 0xF3, 0xE0), ACCENT3)
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.08)
    _p(tf, "\u26A0 Nuance: 77% of fault-triggering tests embed post-hoc developer "
       "knowledge [3]; D4J\u2019s single-fault assumption is artificial \u2014 "
       "real programs average ~9.2 co-existing faults [21]",
       11, RGBColor(0x8B, 0x45, 0x13), first=True)

    # ==================================================================
    # SLIDE 4 — INTRODUCTION: WHY BUG CATEGORIZATION  (SP2)
    # ==================================================================
    slide = new_slide(prs, "Why Bug Categorization Matters in the SDLC", "INTRODUCTION")

    tf = _txt(slide, 0.6, 1.4, 6.0, 5.5)
    _p(tf, "The Need for Structured Categorization", 18, PRIMARY, True,
       after=Pt(8), first=True)
    for b in [
        "• Different bug types behave differently and require different fix strategies [7]",
        "• Categorization is operationally necessary for optimizing debugging [12]",
        "• Automated classification improves triage and maintenance efficiency [8]",
        "• Predicting fault category from bug reports helps prioritize early [9]",
        "• Structured categorization provides actionable in-process measurements [12]",
    ]:
        _p(tf, b, 13, TEXT_DARK, after=Pt(6))

    _p(tf, "", 6)
    _p(tf, "Current Challenge", 18, ACCENT3, True, after=Pt(8))
    for b in [
        "• Human labeling is the gold standard but too expensive at scale [8]",
        "• ML-based accuracy remains imperfect without labeled training data",
        "• This gap motivates the need for an LLM-driven approach",
    ]:
        _p(tf, b, 13, TEXT_DARK, after=Pt(6))

    # right side — SDLC flow
    colors = [PRIMARY, ACCENT3, ACCENT1, PRIMARY, ACCENT2]
    labels = ["Bug Occurs", "Manual Triage", "Categorize Defect",
              "Assign & Fix", "Process Improvement"]
    for i, (lbl, c) in enumerate(zip(labels, colors)):
        y = 1.5 + i * 1.1
        s = _pill(slide, 7.0, y, 5.0, 0.7, c)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.text = lbl; p.font.size = Pt(14); p.font.color.rgb = WHITE
        p.font.name = FONT; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # connector
        if i < 4:
            cn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(9.48), Inches(y + 0.72),
                                        Inches(0.04), Inches(0.35))
            _solid(cn, BORDER); _no_line(cn)

    # bottleneck callout
    s = _box(slide, 9.8, 2.9, 2.5, 1.3,
             RGBColor(0xFF, 0xEB, 0xEE), ACCENT3, 2)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    _p(tf, "\u26A1 BOTTLENECK", 11, ACCENT3, True,
       align=PP_ALIGN.CENTER, first=True)
    _p(tf, "Manual categorization is expensive, inconsistent, and does not scale",
       10, TEXT_DARK, align=PP_ALIGN.CENTER)

    # ==================================================================
    # SLIDE 5 — LIT REVIEW: EXISTING CLASSIFICATIONS  (SP3)
    # ==================================================================
    slide = new_slide(prs, "Existing Defects4J Classification Attempts",
                      "LITERATURE REVIEW")

    cols = [
        ("Structural / Graph-Based", "Van der Spuy & Fischer (2025) [4]",
         ["CFG/DFG-based classification",
          "488 bugs across 7 projects",
          "6 control-flow + 2 data-flow classes",
          "Captures HOW bugs manifest syntactically"], PRIMARY),
        ("Patch-Level Taxonomy", "Sobreira et al. (2018) [2]",
         ["Repair-action classification",
          "395 bugs from D4J v1.2 only",
          "43% involve conditional changes",
          "Manual patch taxonomy extraction"], ACCENT1),
        ("Small-Sample Inspection", "Xuan et al. (2017) [6]",
         ["Random sample of 50 bugs",
          "Identifies fixability patterns",
          "Manual analysis reveals detail",
          "Self-acknowledged validity threat"], ACCENT2),
    ]

    for i, (title, paper, bullets, color) in enumerate(cols):
        x = 0.6 + i * 4.2
        # header
        s = _pill(slide, x, 1.4, 3.8, 0.6, color)
        tf = s.text_frame
        p = tf.paragraphs[0]
        p.text = title; p.font.size = Pt(14); p.font.color.rgb = WHITE
        p.font.name = FONT; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # paper ref
        tf = _txt(slide, x + 0.1, 2.1, 3.6, 0.3)
        _p(tf, paper, 10, TEXT_MED, italic=True, first=True)
        # bullets
        tf = _txt(slide, x + 0.1, 2.5, 3.6, 2.5)
        for j, b in enumerate(bullets):
            _p(tf, f"\u2022 {b}", 12, TEXT_DARK, after=Pt(6),
               first=(j == 0))

    # key observation box
    s = _box(slide, 0.6, 5.5, 12.0, 1.2,
             RGBColor(0xE8, 0xF5, 0xE9), ACCENT2, 2)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.12)
    _p(tf, "Key Observation", 14, PRIMARY, True, after=Pt(4), first=True)
    _p(tf, "None of these approaches apply a semantic, root-cause taxonomy like "
       "ODC to the full Defects4J dataset. All are either structural-only, "
       "patch-level, or limited to small subsets.", 12, TEXT_DARK)

    # ==================================================================
    # SLIDE 6 — LIT REVIEW: LLM-BASED CLASSIFICATION  (SP5)
    # ==================================================================
    slide = new_slide(prs, "LLM-Based Defect Classification: A Viable Path",
                      "LITERATURE REVIEW")

    tf = _txt(slide, 0.6, 1.4, 6.0, 5.0)
    _p(tf, "Evidence Supporting LLM-Driven Classification", 18, PRIMARY, True,
       after=Pt(10), first=True)
    evidence = [
        ("\u2726 Issue Report Classification [10]",
         "LLMs effectively classify issue reports \u2014 validates the LLM-for-defect-classification hypothesis"),
        ("\u2726 Fine-Grained Categorization [11]",
         "LLMs + prompt engineering + few-shot achieve fine-grained bug report categorization"),
        ("\u2726 Chain-of-Thought Prompting [19]",
         "CoT prompting improves LLM classification accuracy by >500% on structured tasks"),
        ("\u2726 ODC-Aligned Root Cause Analysis [18]",
         "LLMs generate structured, ODC-aligned root cause explanations from bug reports & stack traces"),
    ]
    for t, d in evidence:
        _p(tf, t, 13, ACCENT1, True, after=Pt(2))
        _p(tf, f"   {d}", 11, TEXT_DARK, after=Pt(10))

    # comparison box
    s = _box(slide, 7.0, 1.4, 5.6, 5.0)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    _p(tf, "LLMs vs Traditional ML", 16, PRIMARY, True,
       align=PP_ALIGN.CENTER, after=Pt(12), first=True)

    for title, lines, color in [
        ("Traditional ML (SVM/NB)", [
            "~77\u201382% accuracy [13][14]",
            "Requires labeled training data",
            "Sparse text limitations"], ACCENT3),
        ("LLM-Based Approach", [
            "Code-aware reasoning",
            "Zero-shot / few-shot capable",
            "Structured taxonomy guidance"], ACCENT1),
    ]:
        _p(tf, title, 14, color, True, after=Pt(4))
        for ln in lines:
            _p(tf, f"  \u2022 {ln}", 11, TEXT_DARK, after=Pt(2))
        _p(tf, "", 8)

    _p(tf, "Key Advantage: LLMs are code-aware \u2014 they overcome the "
       "text-only limitations of traditional ML [9]", 11, PRIMARY, True)

    # ==================================================================
    # SLIDE 7 — LLM TAXONOMY PROBLEM  (SP6)
    # ==================================================================
    slide = new_slide(prs, "The Taxonomy Problem for LLMs",
                      "LITERATURE REVIEW")

    # problem statement
    tf = _txt(slide, 0.4, 1.35, 12.5, 0.5)
    _p(tf, "LLMs alone are not enough \u2014 without a structured taxonomy, "
       "classification collapses",
       14, ACCENT3, True, first=True)

    problems = [
        ("Hallucination Without Structure",
         "LLMs hallucinate heavily when classifying bugs without "
         "concrete stack traces or a guiding taxonomy [18]",
         ACCENT3),
        ("Inconsistent Labels",
         "Without standardization, different LLM runs produce "
         "different category labels for the same bug [11]",
         PRIMARY),
        ("Incompatible Taxonomies",
         "Ad hoc ML classification yields a proliferation of "
         "incompatible categories that can\u2019t be compared [8]",
         ACCENT1),
        ("Unactionable Results",
         "Classification without a principled framework produces "
         "vague labels with no operational value [14]",
         PURPLE),
    ]

    for i, (title, desc, color) in enumerate(problems):
        col, row = i % 2, i // 2
        x = 0.4 + col * 6.35
        y = 2.1 + row * 1.7
        s = _box(slide, x, y, 6.1, 1.4, LIGHT_BG, color, 2)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        _p(tf, f"\u2716  {title}", 13, color, True, after=Pt(4), first=True)
        _p(tf, desc, 11, TEXT_DARK)

    # transition callout
    s = _box(slide, 0.4, 5.7, 12.5, 1.1,
             RGBColor(0xE8, 0xF5, 0xE9), ACCENT2, 2)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.1)
    _p(tf, "\u2192  Solution: A Guiding Taxonomy", 14, PRIMARY, True,
       first=True, after=Pt(4))
    _p(tf, "A well-established, mathematically rigorous taxonomy is needed to "
       "constrain and guide LLM outputs. That taxonomy is Orthogonal Defect "
       "Classification (ODC) \u2014 a proven framework with 30+ years of "
       "industrial validation [12].",
       11, TEXT_DARK)

    # ==================================================================
    # SLIDE 8 — LIT REVIEW: ODC FRAMEWORK  (SP7)
    # ==================================================================
    slide = new_slide(prs, "Orthogonal Defect Classification (ODC)",
                      "LITERATURE REVIEW")

    tf = _txt(slide, 0.6, 1.4, 12.0, 0.5)
    _p(tf, "Chillarege et al. (1992) [12]: A mathematically rigorous taxonomy "
       "providing actionable, in-process measurements for software development",
       13, TEXT_DARK, italic=True, first=True)

    # opener header
    s = _pill(slide, 0.6, 2.2, 5.5, 0.5, PRIMARY)
    tf = s.text_frame
    _p(tf, "OPENER Attributes (When Defect Found)", 13, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, (attr, desc) in enumerate([
        ("Activity", "Development phase when defect found"),
        ("Trigger", "Testing condition that exposed defect"),
        ("Impact", "Effect on customer / system"),
    ]):
        tf = _txt(slide, 0.8, 2.85 + i * 0.42, 5.0, 0.35)
        _p(tf, f"  \u25B8 {attr}: {desc}", 12, TEXT_DARK, first=True)

    # closer header
    s = _pill(slide, 7.0, 2.2, 5.6, 0.5, ACCENT1)
    tf = s.text_frame
    _p(tf, "CLOSER Attributes (When Defect Fixed)", 13, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, (attr, desc) in enumerate([
        ("Defect Type", "Root cause category (7 types)"),
        ("Target", "Design/Code artifact affected"),
        ("Qualifier", "Missing, Incorrect, or Extraneous"),
        ("Age", "Base, New, or Rewritten code"),
        ("Source", "In-house, Library, or Vendor"),
    ]):
        tf = _txt(slide, 7.2, 2.85 + i * 0.38, 5.4, 0.35)
        _p(tf, f"  \u25B8 {attr}: {desc}", 12, TEXT_DARK, first=True)

    # 7 types header
    s = _pill(slide, 0.6, 4.6, 12.0, 0.5, PRIMARY)
    tf = s.text_frame
    _p(tf, "The 7 ODC Defect Types (Target = Design/Code)", 14, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    left_types = [
        ("Algorithm/Method", "Wrong computational logic", "C&DF"),
        ("Assignment/Initialization", "Wrong value / initialization", "C&DF"),
        ("Checking", "Missing / wrong validation", "C&DF"),
        ("Timing/Serialization", "Concurrency / ordering issue", "C&DF"),
    ]
    right_types = [
        ("Function/Class/Object", "Design-level gap", "Structural"),
        ("Interface/O-O Messages", "Component boundary issue", "Structural"),
        ("Relationship", "Entity association problem", "Structural"),
    ]

    for i, (name, desc, fam) in enumerate(left_types):
        tf = _txt(slide, 0.8, 5.25 + i * 0.38, 5.5, 0.35)
        _p(tf, f"\u25CF {name}  \u2014 {desc}  [{fam}]", 11, TEXT_DARK,
           first=True)

    for i, (name, desc, fam) in enumerate(right_types):
        tf = _txt(slide, 7.0, 5.25 + i * 0.38, 5.5, 0.35)
        _p(tf, f"\u25CF {name}  \u2014 {desc}  [{fam}]", 11, TEXT_DARK,
           first=True)

    # relevance note
    s = _box(slide, 0.6, 6.65, 12.0, 0.45,
             RGBColor(0xE3, 0xF2, 0xFD), ACCENT1)
    tf = s.text_frame; tf.margin_left = Inches(0.2)
    _p(tf, "ODC remains viable in 2024 \u2014 applied to AI-generated code [17]. "
       "Industrial adoption validated by AutoODC [13]. Used by RCEGen for root "
       "cause analysis [18].", 11, PRIMARY, italic=True, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ==================================================================
    # SLIDE 8 — PROBLEM STATEMENT  (SP4)
    # ==================================================================
    slide = new_slide(prs, "Limitations of Existing Classification Approaches",
                      "PROBLEM STATEMENT")

    tf = _txt(slide, 0.6, 1.4, 12.0, 5.5)
    _p(tf, "Critical Limitations in Current Defects4J Classification", 18,
       PRIMARY, True, after=Pt(10), first=True)

    for t, d in [
        ("\u2717  Coverage Gap",
         "No study has fully labeled all 864+ D4J bugs \u2014 existing work covers only 50\u2013488 bugs [2][4][6]"),
        ("\u2717  Structural-Only Taxonomy",
         "CFG/DFG approaches capture HOW bugs manifest, not WHY they exist semantically [4]"),
        ("\u2717  Patch-Level, Not Root-Cause",
         "Repair-action classification describes fix patterns, not the underlying mechanism [2]"),
        ("\u2717  Post-Fix Dependency",
         "Existing ML approaches require labeled training data (D4J has none for ODC) or post-fix code [14]"),
        ("\u2717  Descriptive, Not Diagnostic",
         "Current labels are too low-level or inconsistent for decision-making in triage & QA [7]"),
    ]:
        _p(tf, t, 14, ACCENT3, True, after=Pt(2))
        _p(tf, f"     {d}", 12, TEXT_DARK, after=Pt(12))

    # ==================================================================
    # SLIDE 9 — RESEARCH GAP  (SP4 gap)
    # ==================================================================
    slide = new_slide(prs, "Research Gap", "PROBLEM STATEMENT")

    # gap box
    s = _box(slide, 0.4, 1.4, 12.5, 1.6,
             RGBColor(0xFC, 0xE4, 0xEC), ACCENT3, 2)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.12)
    _p(tf, "The Core Gap", 18, ACCENT3, True, after=Pt(6), first=True)
    _p(tf, "No study has fully labeled all Defects4J bugs. Existing work "
       "covers only 50\u2013488 bugs, uses structural-only taxonomies, or "
       "requires post-fix code \u2014 preventing pre-triage use.",
       13, TEXT_DARK, after=Pt(4))
    _p(tf, "There exists no context-aware, semantically-grounded, scalable "
       "classification of the full Defects4J dataset.",
       13, PRIMARY, True)

    # what's needed — 4 cards
    tf = _txt(slide, 0.4, 3.3, 12.5, 0.4)
    _p(tf, "What\u2019s Needed: A pipeline that is\u2026", 15, PRIMARY, True,
       first=True)

    needs = [
        ("\u2713 Context-Aware", "Multi-modal evidence:\ncode, tests, bug reports,\nstack traces", PRIMARY),
        ("\u2713 Semantic", "Root-cause classification,\nnot syntactic patterns", ACCENT1),
        ("\u2713 Taxonomy-Constrained", "ODC prevents hallucination\n\u0026 ensures comparability", ACCENT2),
        ("\u2713 Scalable", "LLM-based for full dataset\nwithout manual labeling", PURPLE),
    ]
    for i, (label, desc, color) in enumerate(needs):
        x = 0.4 + i * 3.15
        s = _box(slide, x, 3.85, 2.95, 1.7, LIGHT_BG, color, 2)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        _p(tf, label, 13, color, True, after=Pt(5), first=True)
        _p(tf, desc, 10, TEXT_DARK)

    # bridge to next slides
    s = _box(slide, 0.4, 5.85, 12.5, 0.5,
             RGBColor(0xE3, 0xF2, 0xFD), ACCENT1)
    tf = s.text_frame; tf.margin_left = Inches(0.2)
    _p(tf, "\u2192  Our approach: encode Chillarege\u2019s ODC workflow into "
       "an LLM classification pipeline with structured scientific debugging [12][15]",
       11, PRIMARY, True, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ==================================================================
    # SLIDE 10 — OBJECTIVES  (SP8/9/10)
    # ==================================================================
    slide = new_slide(prs, "Research Objectives", "OBJECTIVES")

    objs = [
        ("RO-1", "Build an automated pipeline that classifies Defects4J bugs "
         "into ODC defect types using LLM reasoning guided by structured "
         "scientific debugging", PRIMARY),
        ("RO-2", "Collect bug evidence in two modes \u2014 pre-fix (symptoms only) "
         "and post-fix (with the actual code fix) \u2014 mapping Defects4J "
         "artifacts to ODC opener and closer attributes", ACCENT1),
        ("RO-3", "Apply Zeller\u2019s scientific debugging methodology as an LLM "
         "prompting strategy to produce explainable, step-by-step "
         "classifications", ACCENT2),
        ("RO-4", "Design a 4-tier evaluation framework that compares pre-fix "
         "vs post-fix classifications using strict match, top-2 match, "
         "family match, and Cohen\u2019s Kappa", ACCENT3),
    ]

    for i, (code, text, color) in enumerate(objs):
        y = 1.6 + i * 1.3
        s = _pill(slide, 0.6, y, 0.9, 0.5, color)
        tf = s.text_frame
        _p(tf, code, 13, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        s = _box(slide, 1.7, y, 10.9, 0.85)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
        _p(tf, text, 13, TEXT_DARK, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # analogy box
    s = _box(slide, 0.6, 6.0, 12.0, 0.85,
             RGBColor(0xE3, 0xF2, 0xFD), ACCENT1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "Design Inspiration:", 12, PRIMARY, True, first=True)
    _p(tf, "Just as Kang et al. (2023) [15] encoded Zeller\u2019s scientific "
       "debugging into an LLM pipeline for fault localization, our approach "
       "encodes Chillarege\u2019s ODC workflow into an LLM classification "
       "pipeline with structured prompting.", 12, TEXT_DARK, italic=True)

    # ==================================================================
    # SLIDE 11 — ODC \u00D7 DEFECTS4J MAPPING  (SP8)
    # ==================================================================
    slide = new_slide(prs, "ODC \u00D7 Defects4J: Natural Artifact Mapping",
                      "PROPOSED METHODOLOGY")

    # --- OPENER header ---
    s = _pill(slide, 0.4, 1.4, 6.1, 0.45, PRIMARY)
    tf = s.text_frame
    _p(tf, "OPENER (Discovery Story)", 13, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    opener_rows = [
        ("Tests",         "Activity",  "Test type tells us WHAT was\nbeing done when defect surfaced"),
        ("Stack Traces",  "Trigger",   "Runtime context shows WHAT\ncondition woke the dormant defect"),
        ("Bug Reports",   "Impact",    "User-visible symptom describes\nthe consequence (crash, wrong output)"),
    ]
    for i, (d4j, odc, desc) in enumerate(opener_rows):
        y = 2.05 + i * 1.15
        s = _pill(slide, 0.4, y, 2.2, 0.55, PRIMARY)
        tf = s.text_frame
        _p(tf, d4j, 12, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = _txt(slide, 2.7, y + 0.05, 0.5, 0.45)
        _p(tf, "\u2192", 20, PRIMARY, True, align=PP_ALIGN.CENTER, first=True)
        s = _pill(slide, 3.25, y, 1.6, 0.55, ACCENT1)
        tf = s.text_frame
        _p(tf, odc, 12, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = _txt(slide, 5.0, y, 1.6, 0.85)
        _p(tf, desc, 9, TEXT_DARK, first=True)

    # --- CLOSER header ---
    s = _pill(slide, 6.8, 1.4, 6.1, 0.45, ACCENT1)
    tf = s.text_frame
    _p(tf, "CLOSER (Repair Story)", 13, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    closer_rows = [
        ("Code Diffs",       "Type + Qualifier", "Diff reveals the semantic mistake\n(Type) and how it manifested\n(Missing / Incorrect / Extraneous)"),
        ("Modified Class",   "Target + Age",     "WHERE the fix was applied\n(Target = Code) and WHEN the\ndefect was injected (Base / New)"),
    ]
    for i, (d4j, odc, desc) in enumerate(closer_rows):
        y = 2.05 + i * 1.7
        s = _pill(slide, 6.8, y, 2.2, 0.55, ACCENT1)
        tf = s.text_frame
        _p(tf, d4j, 12, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = _txt(slide, 9.1, y + 0.05, 0.5, 0.45)
        _p(tf, "\u2192", 20, ACCENT1, True, align=PP_ALIGN.CENTER, first=True)
        s = _pill(slide, 9.65, y, 2.0, 0.55, PRIMARY)
        tf = s.text_frame
        _p(tf, odc, 11, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        s = _box(slide, 6.8, y + 0.65, 4.85, 0.85, LIGHT_BG, ACCENT1, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.04)
        _p(tf, desc, 9, TEXT_DARK, first=True)

    # --- Key Insight footer ---
    s = _box(slide, 0.4, 5.65, 12.5, 1.15,
             RGBColor(0xE8, 0xF5, 0xE9), ACCENT2, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "Key Insight:", 11, PRIMARY, True, first=True)
    _p(tf, "Tests \u2192 Stack Traces \u2192 Bug Reports form the Opener chain "
       "(Activity \u2192 Trigger \u2192 Impact) \u2014 the discovery story.", 10, TEXT_DARK)
    _p(tf, "Code Diffs \u2192 Modified Class form the Closer chain "
       "(Type/Qualifier \u2192 Target/Age) \u2014 the repair story.", 10, TEXT_DARK)
    _p(tf, "ODC\u2019s two-phase design keeps subjective discovery data "
       "separate from objective technical data [12].",
       9, TEXT_MED, italic=True)

    # ==================================================================
    # SLIDE 12 — PIPELINE ARCHITECTURE OVERVIEW  (SP10/11)
    # ==================================================================
    slide = new_slide(prs, "Pipeline Architecture Overview", "PROPOSED METHODOLOGY")

    stages = [
        ("1. COLLECT", "Evidence\nGathering",
         "D4J CLI \u2192 multi-modal\nevidence \u2192 context.json", PRIMARY),
        ("2. CLASSIFY", "LLM\nClassification",
         "context.json \u2192 ODC prompt\n\u2192 LLM \u2192 classification.json", ACCENT1),
        ("3. COMPARE", "Evaluation\nFramework",
         "Pre-fix vs Post-fix\n\u2192 multi-tier analysis", ACCENT2),
    ]

    for i, (label, title, desc, color) in enumerate(stages):
        x = 0.6 + i * 4.2
        # card
        _box(slide, x, 1.5, 3.8, 3.0, LIGHT_BG, color, 2)
        # label pill
        s = _pill(slide, x, 1.3, 1.8, 0.4, color)
        tf = s.text_frame
        _p(tf, label, 11, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # title
        tf = _txt(slide, x + 0.2, 1.85, 3.4, 0.7)
        _p(tf, title, 18, color, True, align=PP_ALIGN.CENTER, first=True)
        # desc
        tf = _txt(slide, x + 0.2, 2.65, 3.4, 1.6)
        _p(tf, desc, 12, TEXT_DARK, align=PP_ALIGN.CENTER, first=True)

    # arrows
    for i in range(2):
        tf = _txt(slide, 4.5 + i * 4.2, 2.7, 0.6, 0.5)
        _p(tf, "\u2192", 28, TEXT_LIGHT, True, align=PP_ALIGN.CENTER, first=True)

    # principles
    s = _box(slide, 0.6, 4.8, 12.0, 2.2)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.12)
    _p(tf, "Key Architectural Principles", 14, PRIMARY, True,
       after=Pt(6), first=True)
    for b in [
        "\u25B8 File-oriented, synchronous \u2014 no database, queue, or async orchestration",
        "\u25B8 Multi-provider LLM support (Gemini, OpenRouter, OpenAI-compatible)",
        "\u25B8 Pre-fix evidence by default; post-fix via explicit --include-fix-diff",
        "\u25B8 Hidden oracle exclusion: classes.modified never enters the LLM prompt",
        "\u25B8 Inspired by AutoSD [15]: Zeller\u2019s scientific debugging in structured prompts [16]",
    ]:
        _p(tf, f"  {b}", 11, TEXT_DARK, after=Pt(3))

    # ==================================================================
    # SLIDE 13 — EVIDENCE COLLECTION  (SP11A)
    # ==================================================================
    slide = new_slide(prs, "Evidence Collection: Building context.json",
                      "PROPOSED METHODOLOGY")

    steps = [
        ("1", "Query D4J Metadata",
         "Project, bug info, report URL, modified classes"),
        ("2", "Fetch Bug Report",
         "JIRA / GitHub API / HTML \u2192 opener-side evidence"),
        ("3", "Checkout & Compile",
         "Buggy version checkout \u2192 compile \u2192 test"),
        ("4", "Parse Failures",
         "Extract test names, stack traces, exception headlines"),
        ("5", "Select Suspicious Frames",
         "Filter framework / JDK frames \u2192 project source priority"),
        ("6", "Extract Code Snippets",
         "Production (\u00B112 lines) + test source (\u00B118 lines)"),
        ("7", "Coverage (optional)",
         "Cobertura XML \u2192 line/branch execution rates"),
        ("8", "Serialize context.json",
         "Complete evidence payload for classification"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        col, row = i // 4, i % 4
        x = 0.6 + col * 6.4
        y = 1.5 + row * 1.2

        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.05),
                                   Inches(0.4), Inches(0.4))
        _solid(s, ACCENT1); _no_line(s)
        tf = s.text_frame
        _p(tf, num, 12, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = _txt(slide, x + 0.55, y, 5.5, 0.35)
        _p(tf, title, 13, PRIMARY, True, first=True)
        tf = _txt(slide, x + 0.55, y + 0.35, 5.5, 0.5)
        _p(tf, desc, 11, TEXT_MED, first=True)

    # hidden oracle callout
    s = _box(slide, 0.6, 6.4, 12.0, 0.5,
             RGBColor(0xFF, 0xEB, 0xEE), ACCENT3)
    tf = s.text_frame; tf.margin_left = Inches(0.2)
    _p(tf, "\u26A0 Hidden Oracle: classes.modified stored in hidden_oracles "
       "\u2014 deliberately excluded from LLM prompt to prevent data leakage [3]",
       11, ACCENT3, True, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ==================================================================
    # SLIDE 14 — PRE-FIX VS POST-FIX  (SP11A)
    # ==================================================================
    slide = new_slide(prs, "Evidence Modes: Pre-fix vs Post-fix",
                      "PROPOSED METHODOLOGY")

    # --- Pre-fix card ---
    _box(slide, 0.4, 1.4, 6.1, 4.5, LIGHT_BG, PRIMARY, 2)
    s = _pill(slide, 0.4, 1.2, 6.1, 0.5, PRIMARY)
    tf = s.text_frame
    _p(tf, "PRE-FIX MODE (Default)", 14, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = _txt(slide, 0.7, 1.95, 5.5, 3.8)
    for label, val in [
        ("LLM Sees", "Symptoms only"),
        ("Evidence", "Stack traces, tests, errors,\ncode snippets, bug report"),
        ("Perspective", '"What went wrong?"'),
        ("ODC Side", "Opener (trigger, activity, impact)"),
    ]:
        _p(tf, label, 12, PRIMARY, True, after=Pt(2),
           first=(label == "LLM Sees"))
        _p(tf, val, 11, TEXT_DARK, after=Pt(8))

    # doctor analogy
    s = _box(slide, 0.7, 4.7, 5.5, 0.5)
    tf = s.text_frame; tf.margin_left = Inches(0.1)
    _p(tf, "\u2248  A doctor diagnosing from symptoms", 10, TEXT_MED,
       italic=True, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Post-fix card ---
    _box(slide, 6.8, 1.4, 6.1, 4.5, LIGHT_BG, ACCENT1, 2)
    s = _pill(slide, 6.8, 1.2, 6.1, 0.5, ACCENT1)
    tf = s.text_frame
    _p(tf, "POST-FIX MODE (--include-fix-diff)", 14, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf = _txt(slide, 7.1, 1.95, 5.5, 3.8)
    for label, val in [
        ("LLM Sees", "Symptoms + actual code fix"),
        ("Evidence", "All pre-fix evidence +\nbuggy\u2192fixed diff"),
        ("Perspective", '"What was changed?"'),
        ("ODC Side", "Closer (type, qualifier, age, source)"),
    ]:
        _p(tf, label, 12, ACCENT1, True, after=Pt(2),
           first=(label == "LLM Sees"))
        _p(tf, val, 11, TEXT_DARK, after=Pt(8))

    # pathologist analogy
    s = _box(slide, 7.1, 4.7, 5.5, 0.5)
    tf = s.text_frame; tf.margin_left = Inches(0.1)
    _p(tf, "\u2248  A pathologist with biopsy results", 10, TEXT_MED,
       italic=True, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- VS divider ---
    tf = _txt(slide, 6.1, 3.0, 0.7, 0.6)
    _p(tf, "vs", 18, TEXT_LIGHT, True, align=PP_ALIGN.CENTER, first=True)

    # --- insight bar ---
    s = _box(slide, 0.4, 6.0, 12.5, 0.85,
             RGBColor(0xE8, 0xF5, 0xE9), ACCENT2)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "ODC\u2019s opener/closer model requires evidence from both "
       "perspectives [12]. Pre-fix captures opener evidence, "
       "post-fix captures closer evidence.",
       11, TEXT_DARK, italic=True, first=True)

    # ==================================================================
    # SLIDE 15 — ODC TAXONOMY & PROMPT DESIGN  (SP11B)
    # ==================================================================
    slide = new_slide(prs, "ODC Taxonomy & Prompt Design", "PROPOSED METHODOLOGY")
    tf = _txt(slide, 0.4, 1.35, 12.5, 0.35)
    _p(tf, "System Prompt: 6 Components", 15, PRIMARY, True, first=True)

    prompt_blocks = [
        ("1", "ODC Taxonomy",
         "7 types with definitions,\nindicators \u0026 contrastive\n\u201Cdistinguish_from\u201D guidance",
         PRIMARY),
        ("2", "Anti-Bias Rules",
         "\u201CDo NOT default to\nFunction/Class/Object\u201D \u2014\nsimpler explanations first",
         ACCENT3),
        ("3", "Scientific Debugging",
         "Hypothesize \u2192 Observe \u2192\nPredict \u2192 Examine \u2192\nConclude [15][16]",
         ACCENT1),
        ("4", "Diagnostic Tree",
         "7 questions, each maps\nto one ODC type\n(FCO deliberately last)",
         ACCENT2),
        ("5", "JSON Schema",
         "Structured output:\nreasoning, confidence,\nalternatives, ODC fields",
         PURPLE),
        ("6", "Few-Shot (\u00D75)",
         "Contrastive examples:\nSymptom \u2192 Code \u2192 Type\n\u2192 NOT X [11][19]",
         RGBColor(0x00, 0x89, 0x6B)),
    ]

    for i, (num, title, desc, color) in enumerate(prompt_blocks):
        col, row = i % 3, i // 3
        x = 0.4 + col * 4.2
        y = 1.9 + row * 2.4
        s = _box(slide, x, y, 3.9, 2.1, LIGHT_BG, color, 2)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        _p(tf, f"{num}. {title}", 14, color, True, after=Pt(6), first=True)
        _p(tf, desc, 10, TEXT_DARK)

    # footer
    s = _box(slide, 0.4, 6.5, 12.5, 0.4)
    tf = s.text_frame; tf.margin_left = Inches(0.15)
    _p(tf, "Few-shot pattern: Symptom \u2192 Code \u2192 Classification \u2192 "
       "\u201CNOT X because ...\u201D  |  Anti-bias: forces simpler ODC types "
       "before structural ones [11][19]",
       9, TEXT_MED, italic=True, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ==================================================================
    # SLIDE 16 — SCIENTIFIC DEBUGGING PROTOCOL  (SP11C)
    # ==================================================================
    slide = new_slide(prs, "Scientific Debugging Protocol", "PROPOSED METHODOLOGY")

    step_data = [
        ("HYPOTHESIZE", "Form a root-cause hypothesis\nfrom the bug report", PRIMARY),
        ("OBSERVE", "Examine failure symptoms:\nerror messages, stack traces", ACCENT1),
        ("PREDICT", "What code pattern would\nexplain this defect type?", ACCENT2),
        ("EXAMINE", "Check code snippets\nagainst the prediction", PURPLE),
        ("CONCLUDE", "Select ODC type for the\nroot cause mechanism", ACCENT3),
    ]

    for i, (step, desc, color) in enumerate(step_data):
        x = 0.3 + i * 2.55
        s = _pill(slide, x, 1.5, 2.35, 0.5, color)
        tf = s.text_frame
        _p(tf, f"Step {i+1}: {step}", 11, WHITE, True,
           align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        s = _box(slide, x, 2.1, 2.35, 1.1, LIGHT_BG, color, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.1)
        _p(tf, desc, 10, TEXT_DARK, align=PP_ALIGN.CENTER, first=True)

    # arrows between steps
    for i in range(4):
        tf = _txt(slide, 2.7 + i * 2.55, 1.53, 0.5, 0.45)
        _p(tf, "\u2192", 18, TEXT_LIGHT, True, align=PP_ALIGN.CENTER, first=True)

    # 7-question header
    s = _pill(slide, 0.4, 3.6, 12.5, 0.4, PRIMARY)
    tf = s.text_frame
    _p(tf, "7-Question Diagnostic Decision Tree", 13, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    questions = [
        ("Q1", "Missing guard/condition?", "\u2192 Checking", ACCENT1),
        ("Q2", "Wrong value or init?", "\u2192 Assignment/Init", ACCENT1),
        ("Q3", "Wrong computation logic?", "\u2192 Algorithm/Method", ACCENT1),
        ("Q4", "Component boundary issue?", "\u2192 Interface/O-O", ACCENT1),
        ("Q5", "Execution order dependent?", "\u2192 Timing/Serial", ACCENT2),
        ("Q6", "Entity association issue?", "\u2192 Relationship", ACCENT2),
        ("Q7", "Needs design-level fix?", "\u2192 Fun/Cls/Obj*", ACCENT3),
    ]

    # row 1: Q1-Q4
    for i in range(4):
        q, question, result, clr = questions[i]
        x = 0.4 + i * 3.1
        s = _box(slide, x, 4.2, 2.9, 0.95, LIGHT_BG, clr, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.06)
        _p(tf, f"{q}: {question}", 10, TEXT_DARK, True, first=True, after=Pt(3))
        _p(tf, result, 11, PRIMARY, True)

    # row 2: Q5-Q7
    for i in range(3):
        q, question, result, clr = questions[4 + i]
        x = 0.4 + i * 3.1
        s = _box(slide, x, 5.35, 2.9, 0.95, LIGHT_BG, clr, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.06)
        _p(tf, f"{q}: {question}", 10, TEXT_DARK, True, first=True, after=Pt(3))
        _p(tf, result, 11, PRIMARY, True)

    tf = _txt(slide, 0.4, 6.5, 12.5, 0.35)
    _p(tf, "*Q7 deliberately last \u2014 forces simpler explanations before "
       "\u201Cdesign gap\u201D [19]  |  Based on: Zeller 2009 [16], "
       "Kang et al. 2023 [15]", 9, TEXT_MED, italic=True, first=True)

    # ==================================================================
    # SLIDE 17 — CLASSIFICATION PIPELINE  (SP11D)
    # ==================================================================
    slide = new_slide(prs, "Classification Pipeline & Validation",
                      "PROPOSED METHODOLOGY")

    # --- ROW 1: Prompt Construction ---
    tf = _txt(slide, 0.4, 1.35, 12.5, 0.35)
    _p(tf, "Prompt Construction", 14, PRIMARY, True, first=True)

    row1 = [
        ("context.json", "Load bug\nevidence", PRIMARY),
        ("System Prompt", "ODC taxonomy +\nscientific debug\n+ few-shot", ACCENT1),
        ("User Prompt", "Evidence payload\n(oracles filtered)", ACCENT2),
        ("LLM API Call", "Gemini / OpenRouter\n/ OpenAI-compat", PURPLE),
    ]
    for i, (label, desc, color) in enumerate(row1):
        x = 0.4 + i * 3.15
        s = _pill(slide, x, 1.85, 2.7, 0.55, color)
        tf = s.text_frame
        _p(tf, label, 12, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # description below
        s = _box(slide, x, 2.5, 2.7, 1.05, LIGHT_BG, color, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.06)
        _p(tf, desc, 9, TEXT_DARK, align=PP_ALIGN.CENTER, first=True)
    # arrows between row 1 nodes
    for i in range(3):
        tf = _txt(slide, 3.15 + i * 3.15, 1.88, 0.5, 0.5)
        _p(tf, "\u2192", 20, TEXT_LIGHT, True, align=PP_ALIGN.CENTER, first=True)

    # --- ROW 2: Post-Processing ---
    tf = _txt(slide, 0.4, 3.8, 12.5, 0.35)
    _p(tf, "Post-Processing \u0026 Validation", 14, PRIMARY, True, first=True)

    row2 = [
        ("Extract JSON", "Parse from raw\nLLM response", ACCENT1),
        ("Validate Type", "Check against 7\ncanonical ODC names", ACCENT3),
        ("Canonicalize", "Overwrite family\n+ normalize fields", PRIMARY),
        ("Write Outputs", "classification.json\n+ report.md", ACCENT2),
    ]
    for i, (label, desc, color) in enumerate(row2):
        x = 0.4 + i * 3.15
        s = _pill(slide, x, 4.3, 2.7, 0.55, color)
        tf = s.text_frame
        _p(tf, label, 12, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # description below
        s = _box(slide, x, 4.95, 2.7, 0.85, LIGHT_BG, color, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.06)
        _p(tf, desc, 9, TEXT_DARK, align=PP_ALIGN.CENTER, first=True)
    # arrows between row 2 nodes
    for i in range(3):
        tf = _txt(slide, 3.15 + i * 3.15, 4.33, 0.5, 0.5)
        _p(tf, "\u2192", 20, TEXT_LIGHT, True, align=PP_ALIGN.CENTER, first=True)

    # --- Safeguards footer ---
    s = _box(slide, 0.4, 6.1, 12.5, 0.8,
             RGBColor(0xE3, 0xF2, 0xFD), ACCENT1, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "Safeguards:", 11, PRIMARY, True, first=True)
    _p(tf, "Temperature=0.0 (deterministic)  \u2022  "
       "Hidden oracle exclusion [3]  \u2022  "
       "Family always canonicalized from code  \u2022  "
       "Transient-only retry (429/500)  \u2022  "
       "Schema-enforced JSON (Gemini)",
       9, TEXT_DARK)

    # ==================================================================
    # SLIDE 18 — EVALUATION FRAMEWORK  (SP9/SP13)
    # ==================================================================
    slide = new_slide(prs, "Evaluation Strategy & Framework",
                      "PROPOSED METHODOLOGY")

    tf = _txt(slide, 0.4, 1.35, 12.5, 0.4)
    _p(tf, "Challenge: No ODC ground-truth labels exist for Defects4J "
       "\u2192 we cannot use traditional precision/recall",
       13, ACCENT3, True, first=True)

    tiers = [
        ("Tier 1", "Strict Match", "Do pre-fix and post-fix\nproduce the exact same\nODC defect type?",
         "Expected: 40\u201360%", PRIMARY),
        ("Tier 2", "Top-2 Match", "Does the primary type of\none mode appear in the\nother\u2019s alternatives?",
         "Expected: 60\u201380%", ACCENT1),
        ("Tier 3", "Family Match", "Are both types in the\nsame ODC family?\n(Control&DF / Structural)",
         "Expected: 70\u201390%", ACCENT2),
        ("Tier 4", "Cohen\u2019s Kappa", "Statistical measure of\nagreement beyond chance\n(inter-rater reliability)",
         "Expected: \u03BA \u2265 0.50", PURPLE),
    ]

    for i, (tier, name, desc, exp, color) in enumerate(tiers):
        x = 0.4 + i * 3.15
        s = _box(slide, x, 2.0, 2.95, 3.0, LIGHT_BG, color, 2)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        _p(tf, tier, 11, color, True, after=Pt(2), first=True)
        _p(tf, name, 16, color, True, after=Pt(6))
        _p(tf, desc, 11, TEXT_DARK, after=Pt(8))
        _p(tf, exp, 11, color, True)

    # rationale box
    s = _box(slide, 0.4, 5.3, 12.5, 0.7,
             RGBColor(0xE3, 0xF2, 0xFD), ACCENT1, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "Why this works: ", 11, PRIMARY, True, first=True)
    _p(tf, "Even human ODC experts disagree on 20\u201330% of bugs [12]. "
       "A tiered framework captures both exact agreement and "
       "meaningful near-agreement, matching industry evaluation practice [14].",
       10, TEXT_DARK)

    tf = _txt(slide, 0.4, 6.2, 12.5, 0.3)
    _p(tf, "Based on: Chillarege 1992 [12], Thung 2012 [14]",
       9, TEXT_MED, italic=True, first=True)

    # ==================================================================
    # SLIDE 19 — RESULTS / FINDINGS  (populated with analysis data)
    # ==================================================================
    slide = new_slide(prs, "Results / Findings", "RESULTS")

    # --- Headline stats row ---
    stats = [
        ("34", "Bug Pairs\nAnalyzed", PRIMARY),
        ("17", "Defects4J\nProjects", ACCENT1),
        ("73.5%", "Strict Type\nMatch", ACCENT2),
        ("94.1%", "Top-2\nMatch", PURPLE),
        ("97.1%", "Family\nMatch", ACCENT3),
    ]
    for i, (val, label, color) in enumerate(stats):
        x = 0.4 + i * 2.55
        s = _box(slide, x, 1.35, 2.35, 1.1, LIGHT_BG, color, 2)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.04)
        _p(tf, val, 22, color, True, align=PP_ALIGN.CENTER, first=True)
        _p(tf, label, 9, TEXT_DARK, align=PP_ALIGN.CENTER)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Classification Distribution (type unchanged) ---
    s = _box(slide, 0.4, 2.7, 6.4, 2.3, LIGHT_BG, PRIMARY, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    _p(tf, "Type Agreement Distribution (25/34 stable)", 12, PRIMARY, True,
       after=Pt(5), first=True)
    type_unchanged = [
        ("Algorithm/Method", 16, "\u2588" * 16),
        ("Checking", 6, "\u2588" * 6),
        ("Assignment/Init", 2, "\u2588" * 2),
        ("Function/Class/Obj", 1, "\u2588" * 1),
    ]
    for name, count, bar in type_unchanged:
        _p(tf, f"  {name} ({count})", 9, TEXT_DARK, True, after=Pt(1))
        _p(tf, f"    {bar}", 9, ACCENT1, after=Pt(3))

    # --- Type transitions (type changed) ---
    s = _box(slide, 7.0, 2.7, 5.7, 2.3, LIGHT_BG, ACCENT3, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    _p(tf, "Type Shift Patterns (9/34 changed)", 12, ACCENT3, True,
       after=Pt(4), first=True)
    transitions = [
        ("Checking \u2192 Algorithm/Method", "3"),
        ("Algorithm/Method \u2192 Checking", "2"),
        ("Checking \u2192 Assignment/Init", "1"),
        ("Checking \u2192 Interface/O-O Msg", "1"),
        ("Algorithm/Method \u2192 Assignment/Init", "1"),
        ("Function/Class/Obj \u2192 Relationship", "1"),
    ]
    for trans, cnt in transitions:
        _p(tf, f"  {trans}  [{cnt}]", 9, TEXT_DARK, after=Pt(2))

    # --- Key finding box ---
    s = _box(slide, 0.4, 5.2, 12.3, 0.85,
             RGBColor(0xE8, 0xF5, 0xE9), ACCENT2, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "Key Finding:", 12, PRIMARY, True, first=True)
    _p(tf, "Out of 9 bugs where the type shifted, 7 still had overlap "
       "through alternative types. Only 2 bugs (JacksonCore-24, Math-34) "
       "had no alternative overlap \u2014 but both matched at the family level.",
       10, TEXT_DARK)

    # --- Project coverage bar ---
    s = _box(slide, 0.4, 6.2, 12.3, 0.7)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
    _p(tf, "Projects Covered:  Chart \u2022 Cli \u2022 Closure \u2022 Codec "
       "\u2022 Collections \u2022 Compress \u2022 Csv \u2022 Gson \u2022 JacksonCore "
       "\u2022 JacksonDatabind \u2022 JacksonXml \u2022 Jsoup \u2022 JxPath \u2022 Lang "
       "\u2022 Math \u2022 Mockito \u2022 Time",
       9, TEXT_MED, first=True)

    # ==================================================================
    # SLIDE 20 — DISCUSSION / ANALYSIS  (populated with analysis data)
    # ==================================================================
    slide = new_slide(prs, "Discussion / Analysis", "ANALYSIS")

    # --- Divergence pattern cards (3 across) ---
    divg = [
        ("Exact Match", "25 (73.5%)", "Same type in both modes", ACCENT2),
        ("Soft Divergence", "7 (20.6%)", "Type shifted, cross-alt match",
         RGBColor(0xFF, 0xA7, 0x26)),
        ("Hard Divergence", "2 (5.9%)", "No alternative overlap", ACCENT3),
    ]
    for i, (label, value, desc, clr) in enumerate(divg):
        x = 0.4 + i * 4.2
        s = _box(slide, x, 1.35, 3.9, 1.2, LIGHT_BG, clr, 2)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.06)
        _p(tf, value, 20, clr, True, align=PP_ALIGN.CENTER, first=True)
        _p(tf, label, 12, clr, True, align=PP_ALIGN.CENTER)
        _p(tf, desc, 9, TEXT_DARK, align=PP_ALIGN.CENTER)

    # --- Why diverge (3 compact reasons) ---
    s = _pill(slide, 0.4, 2.8, 12.5, 0.4, PRIMARY)
    tf = s.text_frame
    _p(tf, "Why Do Types Diverge?", 13, WHITE, True,
       align=PP_ALIGN.CENTER, first=True)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    reasons = [
        ("Symptom vs Cause", "Pre-fix sees symptoms;\npost-fix sees the actual fix", ACCENT1),
        ("ODC Boundary Ambiguity", "Checking \u2194 Algorithm boundary\nis inherently ambiguous", ACCENT3),
        ("Information Gap", "Post-fix has code change;\npre-fix infers from signals", PURPLE),
    ]
    for i, (title, desc, color) in enumerate(reasons):
        x = 0.4 + i * 4.2
        s = _box(slide, x, 3.4, 3.9, 1.2, LIGHT_BG, color, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.08)
        _p(tf, title, 11, color, True, first=True, after=Pt(3))
        _p(tf, desc, 9, TEXT_DARK)

    # --- Case study (compact) ---
    s = _box(slide, 0.4, 4.85, 12.5, 1.05,
             RGBColor(0xF3, 0xE5, 0xF5), PURPLE, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.06)
    _p(tf, "Case Study: Closure-140 (Soft Divergence)", 11, PURPLE, True,
       after=Pt(3), first=True)
    _p(tf, "Pre-fix \u2192 Checking (validation/guard issue)  |  "
       "Post-fix \u2192 Algorithm/Method (procedural correction)",
       10, TEXT_DARK, after=Pt(2))
    _p(tf, "\u2192 Cross-alternative match: each primary type appears in the "
       "other side\u2019s alternatives \u2014 exactly what ODC theory predicts [12]",
       9, PRIMARY, True)

    # --- Summary bar ---
    s = _box(slide, 0.4, 6.1, 12.5, 0.75,
             RGBColor(0xE8, 0xF5, 0xE9), ACCENT2, 1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "The 26.5% type-shift rate reflects the expected difference between "
       "symptom-based and cause-based classification \u2014 not pipeline failure. "
       "97.1% family match demonstrates robust semantic consistency [12][14].",
       10, TEXT_DARK, first=True)

    # ==================================================================
    # SLIDE 21 — LIMITATIONS  (SP14)
    # ==================================================================
    slide = new_slide(prs, "Limitations", "LIMITATIONS")

    lims = [
        ("L1: No Ground-Truth Labels",
         "Defects4J has no ODC labels from human experts. "
         "We evaluate using self-consistency, not a gold standard. [8][12]"),
        ("L2: LLM Reproducibility",
         "Outputs can vary between models and runs. "
         "Temperature=0.0 helps but doesn\u2019t fully eliminate variability. [11]"),
        ("L3: Data Contamination",
         "LLMs may have seen D4J bugs during training, "
         "especially popular projects like Lang and Math. [3]"),
        ("L4: Heuristic Hints",
         "Opener/closer hints are keyword-based, not validated "
         "by trained assessors. [12]"),
        ("L5: Evidence Quality",
         "Not all projects give equally rich evidence \u2014 "
         "stack trace depth varies between projects."),
        ("L6: Java-Only Scope",
         "Evidence collection is tied to D4J\u2019s Java ecosystem. "
         "Classification logic itself is language-independent."),
    ]

    for i, (title, desc) in enumerate(lims):
        col, row = i // 3, i % 3
        x = 0.4 + col * 6.35
        y = 1.5 + row * 1.7
        clr = ACCENT3 if col == 0 else ACCENT1
        s = _box(slide, x, y, 6.1, 1.4, LIGHT_BG, clr, 1)
        tf = s.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        _p(tf, title, 12, clr, True, after=Pt(3), first=True)
        _p(tf, desc, 10, TEXT_DARK)

    # ==================================================================
    # SLIDE 22 — FUTURE WORK  (SP15)
    # ==================================================================
    slide = new_slide(prs, "Future Work", "FUTURE WORK")

    fw = [
        ("FW-1", "Full Dataset Classification",
         "Scale to all 864+ active D4J bugs \u2014 create the first "
         "complete ODC-labeled Defects4J dataset [9]", PRIMARY),
        ("FW-2", "Multi-Model Consensus",
         "Classify with 3+ models (Gemini, GPT-4, Claude) and "
         "use majority voting to reduce variability [10]", ACCENT1),
        ("FW-3", "Validated ODC Opener",
         "Replace keyword heuristics with proper Activity \u2192 Trigger "
         "mapping using D4J test metadata [12]", ACCENT2),
        ("FW-4", "ODC-Guided APR",
         "Use ODC labels to pick the right repair tool: "
         "e.g. Checking \u2192 conditional templates [5]", PURPLE),
        ("FW-5", "Cross-Dataset Extension",
         "Adapt the pipeline for other datasets like Bears, "
         "BugSwarm, and non-Java projects [7]", ACCENT3),
    ]

    for i, (code, title, desc, color) in enumerate(fw):
        y = 1.5 + i * 1.0
        s = _pill(slide, 0.6, y, 0.8, 0.45, color)
        tf = s.text_frame
        _p(tf, code, 10, WHITE, True, align=PP_ALIGN.CENTER, first=True)
        s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = _txt(slide, 1.6, y - 0.05, 11.0, 0.3)
        _p(tf, title, 14, color, True, first=True)
        tf = _txt(slide, 1.6, y + 0.28, 11.0, 0.5)
        _p(tf, desc, 11, TEXT_DARK, first=True)

    # roadmap bar
    s = _box(slide, 0.4, 6.3, 12.5, 0.7,
             RGBColor(0xE3, 0xF2, 0xFD), ACCENT1)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.06)
    _p(tf, "Next Semester: Full-dataset classification (FW-1) \u2192 "
       "multi-model consensus (FW-2) \u2192 validated opener (FW-3) "
       "\u2192 thesis defense with complete results",
       11, PRIMARY, True, first=True)

    # ==================================================================
    # SLIDE 23 — CONCLUSION  (SP16)
    # ==================================================================
    slide = new_slide(prs, "Conclusion", "CONCLUSION")

    conclusions = [
        ("Research Gap Addressed",
         "First study to apply ODC taxonomy to Defects4J "
         "using LLM-driven scientific debugging"),
        ("Pipeline Contribution",
         "Fully automated pipeline: evidence collection "
         "\u2192 structured prompting \u2192 explainable classifications"),
        ("Methodological Innovation",
         "Combined Chillarege\u2019s ODC with Zeller\u2019s scientific debugging "
         "into a principled LLM prompting strategy"),
        ("4-Tier Evaluation",
         "Strict match, top-2 match, family match, and Cohen\u2019s Kappa "
         "\u2014 accounts for the 20\u201330% human disagreement ceiling [12]"),
        ("Practical Impact",
         "Pre-fix ODC classification works \u2014 type shifts between "
         "modes are expected, not failures"),
    ]

    for i, (title, desc) in enumerate(conclusions):
        y = 1.4 + i * 0.9
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.6), Inches(y), Inches(0.08), Inches(0.6))
        _solid(s, ACCENT1); _no_line(s)

        tf = _txt(slide, 0.9, y, 11.5, 0.3)
        _p(tf, f"\u2726  {title}", 14, PRIMARY, True, first=True)
        tf = _txt(slide, 1.2, y + 0.3, 11.2, 0.4)
        _p(tf, desc, 12, TEXT_DARK, first=True)

    # closing statement
    s = _pill(slide, 0.6, 6.0, 12.0, 0.9, PRIMARY)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.06)
    _p(tf, "\u201CLLM-driven scientific debugging, constrained by ODC taxonomy, "
       "can produce explainable and defensible bug classifications "
       "from pre-fix evidence alone.\u201D",
       12, WHITE, italic=True, align=PP_ALIGN.CENTER, first=True)

    # ==================================================================
    # SLIDE 24 — THANK YOU
    # ==================================================================
    _page[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = PRIMARY
    _page_num(slide, _page[0])

    tf = _txt(slide, 1, 2.5, 11.333, 1.5)
    _p(tf, "Thank You", 48, WHITE, True, align=PP_ALIGN.CENTER, first=True)
    _p(tf, "Questions & Discussion", 24, ACCENT2,
       align=PP_ALIGN.CENTER)

    tf = _txt(slide, 2, 5.0, 9.333, 1.5)
    _p(tf, "Hasin Mahtab Alvee  \u2022  Mohammad Nahiyan Kabir  "
       "\u2022  Md. Sakib Hossain", 14, RGBColor(0xA0, 0xC4, 0xD0),
       align=PP_ALIGN.CENTER, after=Pt(8), first=True)
    _p(tf, "Supervised by: Lutfun Nahar Lota (Asst. Prof.) & "
       "Ishmam Tashdeed (Lecturer)", 12, RGBColor(0xA0, 0xC4, 0xD0),
       align=PP_ALIGN.CENTER, after=Pt(4))
    _p(tf, "Network and Data Analysis Group (NDAG) | "
       "Signal and Speech Processing Laboratory (SSL)", 11, TEXT_LIGHT,
       align=PP_ALIGN.CENTER, after=Pt(2))
    _p(tf, "Islamic University of Technology", 11, TEXT_LIGHT,
       align=PP_ALIGN.CENTER)

    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35),
                               SLIDE_WIDTH, Inches(0.15))
    _solid(s, ACCENT1); _no_line(s)

    # ==================================================================
    # REFERENCE SLIDES (not counted in 24-page limit)
    # ==================================================================
    def ref_slide(title, refs):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
        _slide_title(sl, title, "REFERENCES")
        _bot_bar(sl)
        tf = _txt(sl, 0.6, 1.4, 12.0, 5.5)
        for i, r in enumerate(refs):
            _p(tf, r, 10, TEXT_DARK, after=Pt(7), first=(i == 0))
        return sl

    ref_slide("References (1/2)", [
        "[1] Just, R., Jalali, D., & Ernst, M. D. (2014). Defects4J: A database "
        "of existing faults to enable controlled testing studies for Java "
        "programs. ISSTA 2014, ACM, pp. 437\u2013440.",

        "[2] Sobreira, V., Durieux, T., Madeiral, F., Monperrus, M., & Maia, "
        "M. A. (2018). Dissection of a Bug Dataset: Anatomy of 395 Patches from "
        "Defects4J. IEEE SANER 2018, pp. 130\u2013140.",

        "[3] Rafi, M. N., Chen, A. R., Chen, T.-H. P., & Wang, S. (2024). "
        "Revisiting Defects4J for Fault Localization in Diverse Development "
        "Scenarios. arXiv:2402.13040.",

        "[4] Van der Spuy, A., & Fischer, B. (2025). An Anatomy of 488 Faults "
        "from Defects4J Based on the Control- and Data-Flow Graph "
        "Representations of Programs. arXiv:2502.02299.",

        "[5] Durieux, T., Martinez, M., Monperrus, M., & Wuttke, J. (2018). "
        "Automatic Repair of Real Bugs in Java: A Large-Scale Experiment on the "
        "Defects4J Dataset. arXiv:1811.02429.",

        "[6] Xuan, J., Monperrus, M., et al. (2017). Better test cases for "
        "better automated program repair: A manual inspection of Defects4J bugs. "
        "FSE 2017, ACM, pp. 831\u2013841.",

        "[7] Catolino, G., Palomba, F., Zaidman, A., & Ferrucci, F. (2019). "
        "Not all bugs are the same: Understanding, characterizing, and "
        "classifying bug types. J. Systems and Software, 153, pp. 1\u201318.",

        "[8] Andrade, R., Teixeira, C., Laranjeiro, N., & Vieira, M. (2025). "
        "An Empirical Study on the Classification of Bug Reports with Machine "
        "Learning. arXiv:2503.00660.",

        "[9] Hirsch, T., & Hofer, B. (2022). Using textual bug reports to "
        "predict the fault category of software bugs. Array, 15, 100189.",

        "[10] Colavito, G., Lanubile, F., Novielli, N., & Quaranta, L. (2024). "
        "Large Language Models for Issue Report Classification. Ital-IA 2024, "
        "CEUR Workshop Proceedings.",

        "[11] Koyuncu, A. (2025). Exploring fine-grained bug report "
        "categorization with LLMs and prompt engineering. ACM TOSEM.",
    ])

    ref_slide("References (2/2)", [
        "[12] Chillarege, R., et al. (1992). Orthogonal Defect Classification "
        "\u2014 A Concept for In-Process Measurements. IEEE TSE, 18(11), "
        "pp. 943\u2013956.",

        "[13] Huang, B., & Ng, V., et al. (2011). AutoODC: Automated generation "
        "of Orthogonal Defect Classifications. ASE 2011.",

        "[14] Thung, F., Lo, D., et al. (2012). Automatic Defect Categorization. "
        "WCRE 2012.",

        "[15] Kang, S., Chen, B., Yoo, S., & Lou, J.-G. (2023). Explainable "
        "Automated Debugging via Large Language Model-driven Scientific "
        "Debugging. arXiv:2304.02195.",

        "[16] Zeller, A. (2009). Why Programs Fail: A Guide to Systematic "
        "Debugging (2nd ed.). Morgan Kaufmann.",

        "[17] Pan, X., et al. (2024). Understanding Defects in Generated Codes "
        "by Language Models. CASCON 2024. (arXiv:2408.13372)",

        "[18] Gao, X., et al. (2024). RCEGen: A Generative Approach for "
        "Automated Root Cause Analysis. MDPI Computers, 4(4), 29.",

        "[19] Nong, Y., et al. (2024). Chain-of-Thought Prompting of Large "
        "Language Models for Discovering and Fixing Software Vulnerabilities. "
        "arXiv:2402.17230.",

        "[20] IBM. (2013). Orthogonal Defect Classification v 5.2 for Software "
        "Design and Code. IBM Research.",

        "[21] Callaghan, D. (2024). Mining Bug Repositories for Multi-Fault "
        "Programs. defects4j-mf, GitHub.",
    ])

    # ==================================================================
    # SAVE
    # ==================================================================
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "Thesis_PreDefense_Presentation.pptx")
    prs.save(out)
    print(f"\n[OK] Presentation saved to: {out}")
    print(f"     Total slides: {len(prs.slides)} (25 content + 2 reference)")
    return out


if __name__ == "__main__":
    create_presentation()
